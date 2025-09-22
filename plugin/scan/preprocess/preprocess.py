import os
from typing import Callable
import zipfile
import rarfile
import re
import shutil
import tempfile
import comtypes.client
import time
import pythoncom
import win32com.client as win32
import magic
import fitz
from PIL import Image
import pytesseract
from docx import Document
from docx.oxml import parse_xml
from docx.shared import Pt
import win32con
import win32gui
import win32api
from openpyxl import load_workbook
import xlrd
from pptx import Presentation
import csv
from plugin.config.config import cfg

# 配置 OCR 工具路径
pytesseract.pytesseract.tesseract_cmd = shutil.which('tesseract.exe')#查找tesseract_cmd工具

class OfficeSecurityConfig:
    """Office 安全配置管理器"""
    def __enter__(self):
        self.original_values = {}
        try:
            self.shell = win32.Dispatch("WScript.Shell")
            versions = ['16.0', '15.0', '14.0', '12.0']
            for ver in versions:
                self._modify_registry(f"HKEY_CURRENT_USER\\Software\\Microsoft\\Office\\{ver}\\Word\\Security", [
                    ("VBAWarnings", 1),
                    ("AccessVBOM", 1),
                    ("BlockContentExecutionFromInternet", 0)
                ])
                self._modify_registry(f"HKEY_CURRENT_USER\\Software\\Microsoft\\Office\\{ver}\\Excel\\Security", [
                    ("VBAWarnings", 1),
                    ("AccessVBOM", 1)
                ])
                self._modify_registry(f"HKEY_CURRENT_USER\\Software\\Microsoft\\Office\\{ver}\\PowerPoint\\Security", [
                    ("VBAWarnings", 1),
                    ("AccessVBOM", 1)
                ])
        except Exception as e:
            print(f"安全配置警告: {str(e)}")
        return self

    def _modify_registry(self, key_path, values):
        for name, value in values:
            full_key = f"{key_path}\\{name}"
            try:
                self.original_values[full_key] = self.shell.RegRead(full_key)
            except:
                self.shell.RegWrite(full_key, value, "REG_DWORD")
            else:
                self.shell.RegWrite(full_key, value, "REG_DWORD")

    def __exit__(self, exc_type, exc_val, exc_tb):
        pass

def clear_com_cache():
    cache_dir = os.path.join( os.environ['Temp'],'gen_py')
    if os.path.exists(cache_dir):
        shutil.rmtree(cache_dir, ignore_errors=True)

def detect_file_type(file_path):
    """改进型文件类型检测"""
    try:
        mime = magic.Magic(mime=True)
        mime_type = mime.from_file(file_path)

        type_map = {
            'application/pdf': 'PDF',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'DOCX',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'XLSX',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'PPTX',
            'application/msword': 'DOC',
            'application/vnd.ms-excel': 'XLS',
            'application/vnd.ms-powerpoint': 'PPT',
            'application/zip': 'ZIP',
            'application/x-rar-compressed': 'RAR',
            'inode/x-empty': 'Empty',
            'text/rtf': 'RTF',
            'text/plain': 'Text',
            'text/markdown': 'Markdown',
            'image/png': 'Image',
            'image/jpeg': 'Image',
            'image/gif': 'Image',
            'image/tiff': 'Image',
            'image/bmp': 'Image',
            'application/octet-stream': 'Unknown'
        }

        file_type = type_map.get(mime_type, 'Unknown')

        # 深度检测Office文档
        if file_type in ['ZIP', 'Unknown']:
            try:
                with zipfile.ZipFile(file_path, 'r') as zf:
                    if 'word/document.xml' in zf.namelist():
                        return 'DOCX'
                    elif 'xl/workbook.xml' in zf.namelist():
                        return 'XLSX'
                    elif 'ppt/presentation.xml' in zf.namelist():
                        return 'PPTX'
            except (zipfile.BadZipFile, KeyError):
                pass

        # 检测旧版Office文档
        if file_type in ['Unknown']:
            with open(file_path, 'rb') as f:
                header = f.read(8)
                if header.startswith(b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1'):
                    content = f.read(1024)
                    if b'WordDocument' in content:
                        return 'DOC'
                    elif b'Workbook' in content:
                        return 'XLS'
                    elif b'PowerPoint' in content:
                        return 'PPT'

        # 检测Markdown
        if file_type == 'Text':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                sample = f.read(1024)
                if re.search(r'^#+\s+|!\[.*?\]\(.*?\)|\[.*?\]:', sample):
                    return 'Markdown'

        return file_type
    except Exception as e:
        print(f"文件类型检测失败: {str(e)}")
        return 'Unknown'
def get_handler(file_type, temp_dir):
    """根据文件类型获取对应的处理函数"""
    return {
        'DOCX': lambda path: process_word(path),
        'DOC': lambda path: process_word_legacy(path),
        'XLSX': lambda path: process_excel(path),
        'XLS': lambda path: process_excel_legacy(path),
        'PPTX': lambda path: process_pptx(path),
        'PPT': lambda path: process_ppt_legacy(path),
        'PDF': lambda path: process_pdf(path, temp_dir),
        'Image': lambda path: extract_text_from_image(path),
        'Text': lambda path: open(path, 'r', encoding='utf-8', errors='ignore').read(),
        'Markdown': lambda path: open(path, 'r', encoding='utf-8').read(),
        'RTF': lambda path: open(path, 'r', encoding='utf-8', errors='ignore').read(),
        'Unknown': lambda path: process_unknown(path),
        'Empty': lambda path: process_empty(path),
    }.get(file_type, None)
def rename_file_based_on_type(file_path,outdir):
    """根据文件类型重命名文件"""
    file_type = detect_file_type(file_path)
    file_dir, file_name = os.path.split(file_path)
    base_name, _ = os.path.splitext(file_name)

    extension_map = {
        'DOCX': '.docx',
        'DOC': '.doc',
        'XLSX': '.xlsx',
        'XLS': '.xls',
        'PPTX': '.pptx',
        'PPT': '.ppt',
        'PDF': '.pdf',
        'Image': '.png',
        'Text': '.txt',
        'Markdown': '.md',
        'RTF': '.rtf',
        'ZIP': '.zip',
        'RAR': '.rar',
        'Empty': '.empty',
        'Unknown': '.bin'
    }

    new_extension = extension_map.get(file_type, '.unknown')
    new_file_name = f"{base_name}{new_extension}"
    new_file_path = os.path.join(outdir, new_file_name)
    shutil.copyfile(file_path,new_file_path)
    # os.rename(file_path, new_file_path)
    return new_file_path

def process_compressed_content(file_path, temp_dir):
    """处理压缩文件内容并返回合并的文本"""
    all_content = []
    try:
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                for file_info in zip_ref.infolist():
                    if file_info.is_dir():
                        continue
                    # 处理文件名编码
                    try:
                        file_name = file_info.filename.encode('cp437').decode('utf-8')
                    except UnicodeDecodeError as e:
                        file_name = file_info.filename
                    except UnicodeEncodeError as e:
                        file_name = file_info.filename
                    
                    # 创建临时目录处理嵌套文件
                    with tempfile.TemporaryDirectory(dir=temp_dir) as file_temp_dir:
                        # 提取文件内容到临时文件
                        with zip_ref.open(file_info) as f:
                            content = f.read()
                            tmp_path = os.path.join(file_temp_dir, os.path.basename(file_name))
                            with open(tmp_path, 'wb') as tmp_file:
                                tmp_file.write(content)
                            
                            # 处理文件内容
                            file_type = detect_file_type(tmp_path)
                            handler = get_handler(file_type, file_temp_dir)
                            processed = handler(tmp_path) if handler else ""
                            all_content.append(f"{processed}\n")
        
        elif rarfile.is_rarfile(file_path):
            with rarfile.RarFile(file_path, 'r') as rar_ref:
                for file_info in rar_ref.infolist():
                    if file_info.is_dir():
                        continue
                    file_name = file_info.filename
                    
                    with tempfile.TemporaryDirectory(dir=temp_dir) as file_temp_dir:
                        # 提取文件内容到临时文件
                        with rar_ref.open(file_info) as f:
                            content = f.read()
                            tmp_path = os.path.join(file_temp_dir, os.path.basename(file_name))
                            with open(tmp_path, 'wb') as tmp_file:
                                tmp_file.write(content)
                            
                            # 处理文件内容
                            file_type = detect_file_type(tmp_path)
                            handler = get_handler(file_type, file_temp_dir)
                            processed = handler(tmp_path) if handler else ""
                            all_content.append(f"=== {file_name} ===\n{processed}\n")
    
    except Exception as e:
        print(f"处理压缩文件内容失败: {str(e)}")
    return all_content

def process_excel_legacy(input_path):
    try:
        wb = xlrd.open_workbook(input_path)
        sheet = wb.sheet_by_index(0)
        text = ""
        for row_idx in range(sheet.nrows):
            text += "\t".join([str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]) + "\n"
        return text
    except Exception as e:
        print(f"处理 XLS 文件 {input_path} 时发生错误: {e}")
        return ""

def process_word(input_path, temp_dir=None):
    """直接解析 DOCX 的 XML 内容"""
    from xml.etree.ElementTree import XML
    text = []

    with zipfile.ZipFile(input_path) as z:
        # 提取主文档内容
        doc_content = z.read('word/document.xml')
        tree = XML(doc_content)
        for elem in tree.iter():
            if elem.tag.endswith('t'):
                text.append(elem.text.strip() if elem.text else '')

        # 提取脚注
        if 'word/footnotes.xml' in z.namelist():
            footnotes = z.read('word/footnotes.xml')
            tree = XML(footnotes)
            text.append("\n")
            for elem in tree.iter():
                if elem.tag.endswith('t'):
                    text.append(elem.text.strip() if elem.text else '')

        # 提取批注
        if 'word/comments.xml' in z.namelist():
            comments = z.read('word/comments.xml')
            tree = XML(comments)
            text.append("\n")
            for elem in tree.iter():
                if elem.tag.endswith('t'):
                    text.append(elem.text.strip() if elem.text else '')

    return "\n".join(filter(None, text))

def process_word_legacy(input_path, temp_dir=None):
    """处理 DOC 文件（需要 Windows 环境）"""
    try:
        word = win32.Dispatch("Word.Application")
        doc = word.Documents.Open(input_path)
        text = doc.Content.Text
        doc.Close()
        word.Quit()
        return text
    except Exception as e:
        print(f"处理 DOC 文件 {input_path} 时发生错误: {e}")
        return ""

def process_excel(file_path):
    """增强型Excel处理"""
    try:
        file_type = detect_file_type(file_path)

        if file_type == 'XLS':
            try:
                wb = xlrd.open_workbook(file_path)
                text = []
                for sheet in wb.sheets():
                    for row_idx in range(sheet.nrows):
                        text.append(", ".join(
                            str(xlrd.xldate.xldate_as_datetime(cell, wb.datemode))
                            if cell_type == xlrd.XL_CELL_DATE
                            else str(cell)
                            for cell, cell_type in zip(sheet.row_values(row_idx), sheet.row_types(row_idx))
                        ))
                return "\n".join(text)
            except Exception as e:
                print(f"旧版Excel处理失败: {str(e)}")
                return ""

        try:
            wb = load_workbook(file_path, read_only=True, data_only=True)
        except:
            try:
                from openpyxl import repair
                wb = repair.workbook(file_path)
            except Exception as e:
                print(f"Excel文件修复失败: {str(e)}")
                return ""

        text = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if isinstance(cell, str):
                        row_text.append(cell.strip())
                    elif cell is not None:
                        row_text.append(str(cell))
                text.append(", ".join(row_text))
        return "\n".join(filter(None, text))
    except Exception as e:
        print(f"Excel处理失败: {str(e)}")
        return ""

def process_pptx(file_path):
    """处理新版PPTX文件"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.strip())
                elif shape.has_table:
                    for row in shape.table.rows:
                        text += [cell.text_frame.text.strip() for cell in row.cells]
        return "\n".join(filter(None, text))
    except Exception as e:
        print(f"PPTX处理失败: {str(e)}")
        return ""

def process_ppt_legacy(input_path, temp_dir=None):
    """使用 comtypes 处理 PPT 文件 (Windows + 已激活 Office)"""
    pythoncom.CoInitialize()
    powerpoint = None
    presentation = None

    try:
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        hwnd = win32gui.FindWindow(None, "Microsoft PowerPoint")
        win32gui.ShowWindow(hwnd, win32con.SW_HIDE)
        powerpoint.DisplayAlerts = False

        start_time = time.time()
        presentation = powerpoint.Presentations.Open(
            os.path.abspath(input_path),
            WithWindow=False
        )

        text = []
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    if shape.TextFrame.HasText:
                        content = shape.TextFrame.TextRange.Text.strip()
                        if content:
                            text.append(_clean_text(content))

                if shape.HasTable:
                    table = shape.Table
                    for row in range(1, table.Rows.Count + 1):
                        row_data = []
                        for col in range(1, table.Columns.Count + 1):
                            cell = table.Cell(row, col)
                            cell_text = cell.Shape.TextFrame.TextRange.Text.strip()
                            row_data.append(cell_text)
                        text.append(" | ".join(row_data))

                if shape.Type == 20:
                    text.append(_process_smartart(shape))

        return "\n".join(text)
    except pythoncom.com_error as e:
        error_msg = f"COM 错误 (HRESULT: {e.hresult:#x}): {e.excepinfo[2] if e.excepinfo else '未知错误'}"
        print(f"处理失败: {error_msg}")
        return ""
    except Exception as e:
        print(f"未知错误: {str(e)}")
        return ""
    finally:
        try:
            if presentation:
                presentation.Close()
            if powerpoint:
                powerpoint.Quit()
            os.system('taskkill /f /im POWERPNT.EXE >nul 2>&1')
        except:
            pass
        pythoncom.CoUninitialize()

def _clean_text(text):
    text = re.sub(r'\x0b', '', text)
    text = re.sub(r'\r\n', '\n', text)
    return text.strip()

def _process_smartart(shape):
    content = []
    try:
        nodes = shape.SmartArt.AllNodes
        for node in nodes:
            if node.TextFrame2.HasText:
                content.append(node.TextFrame2.TextRange.Text.strip())
            if node.Nodes.Count > 0:
                content.append(_process_smartart_node(node))
    except:
        pass
    return "\n".join(content)

def _process_smartart_node(node):
    sub_content = []
    for sub_node in node.Nodes:
        if sub_node.TextFrame2.HasText:
            sub_content.append(sub_node.TextFrame2.TextRange.Text.strip())
        if sub_node.Nodes.Count > 0:
            sub_content.append(_process_smartart_node(sub_node))
    return " -> ".join(sub_content)

def process_pdf(file_path, temp_dir):
    """增强型PDF处理"""
    try:
        doc = fitz.open(file_path)
        text = []

        for page in doc:
            text.append(page.get_text())

        img_dir = os.path.join(temp_dir, "pdf_images")
        os.makedirs(img_dir, exist_ok=True)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_img = doc.extract_image(xref)
                img_path = os.path.join(img_dir, f"page{page_num}_img{img_index}.{base_img['ext']}")
                with open(img_path, "wb") as f:
                    f.write(base_img["image"])
                text.append(f"[图片内容：{extract_text_from_image(img_path)}]")

        return "\n".join(text)
    except Exception as e:
        print(f"PDF处理失败: {str(e)}")
        return ""

def extract_text_from_image(image_path):
    """增强型OCR处理"""
    try:
        img = Image.open(image_path)
        return pytesseract.image_to_string(
            img,
            lang='chi_sim+eng',
            config='--psm 3 --oem 1'
        )
    except Exception as e:
        print(f"图片识别失败: {str(e)}")
        return ""

def process_unknown(file_path):
    """未知文件处理"""
    try:
        with open(file_path, 'rb') as f:
            content = f.read(1048576)
            try:
                return content.decode('utf-8', errors='replace')
            except UnicodeDecodeError:
                return "二进制文件内容（不可读）"
    except Exception as e:
        print(f"未知文件处理失败: {str(e)}")
        return ""

def process_empty(file_path):
    """空文件处理"""
    return "空文件"

def clean_content(content):
    if isinstance(content,list):
        cleaned=[]
        for i in content:
            if not i.strip():
                return "识别失败"
            i=re.sub(r'\s+', ' ', i)
            i= re.sub(r'[\x00-\x1F\x7F]', '', i)
            cleaned.append(i.strip())
        return cleaned
    """清理文本内容"""
    if not content.strip():
        return "识别失败"
    cleaned = re.sub(r'\s+', ' ', content)
    cleaned = re.sub(r'[\x00-\x1F\x7F]', '', cleaned)
    cleaned = cleaned.strip()
    return cleaned if cleaned else "识别失败"

def split_content(content, max_length=32767):
    """分割内容到多个列"""
    if content == "识别失败":
        return [content]
    parts = []
    while len(content) > 0:
        part = content[:max_length]
        parts.append(part)
        content = content[max_length:]
    return parts
#TODO 递归处理文件夹下的所有文件和文件夹
def preprocess_recursion(input_path, temp_dir,ScanCallback:Callable[[str,str,str,any],None], processed_files=None):
    if processed_files is None:
        processed_files = set()
    if os.path.isfile(input_path):
        file_id = f"{os.path.abspath(input_path)}-{os.path.getmtime(input_path)}"
        if file_id in processed_files:
            return
        processed_files.add(file_id)

        input_path = rename_file_based_on_type(input_path,temp_dir)
        file_type = detect_file_type(input_path)
        print(f"正在处理: {input_path} [类型: {file_type}]")
        handlers = {
            'DOCX': lambda: process_word(input_path, temp_dir),
            'DOC': lambda: process_word_legacy(input_path),
            'XLSX': lambda: process_excel(input_path),
            'XLS': lambda: process_excel_legacy(input_path),
            'PPTX': lambda: process_pptx(input_path),
            'PPT': lambda: process_ppt_legacy(input_path),
            'PDF': lambda: process_pdf(input_path, temp_dir),
            'Image': lambda: extract_text_from_image(input_path),
            'Text': lambda: open(input_path, 'r', encoding='utf-8', errors='ignore').read(),
            'Markdown': lambda: open(input_path, 'r', encoding='utf-8').read(),
            'ZIP': lambda: process_compressed(input_path, os.path.join(temp_dir, "zip_extract")),
            'RAR': lambda: process_compressed(input_path, os.path.join(temp_dir, "rar_extract")),
            'Unknown': lambda: process_unknown(input_path),
            'Empty': lambda: process_empty(input_path),
            'RTF': lambda: open(input_path, 'r', encoding='utf-8', errors='ignore').read(),
        }

        result = handlers.get(file_type, lambda: "")()

        # 处理压缩文件
        if file_type in ['ZIP', 'RAR']:
            for extracted_file in result:
                preprocess_zip(extracted_file, temp_dir,ScanCallback, processed_files)
            return

        # 处理文件名
        filename = os.path.basename(input_path)
        base_name = os.path.splitext(filename)[0]
        clean_name = re.sub(r'[^A-Za-z0-9]', '', base_name)

        # 处理内容
        cleaned_content = clean_content(result)
        content_parts = split_content(cleaned_content)
        ScanCallback(input_path,clean_name,file_type,content_parts)
        # 写入CSV

    elif os.path.isdir(input_path):
        for item in os.listdir(input_path):
            preprocess_recursion(os.path.join(input_path, item), temp_dir,ScanCallback, processed_files)
#线性处理文件夹下的文件(不包含文件夹)
def preprocess_line(scandir, temp_dir,signals, processed_files=None):
    for item in os.listdir(scandir):
        fullpath=os.path.join(scandir,item)
        if os.path.isdir(fullpath):
            continue
        input_path = rename_file_based_on_type(fullpath,temp_dir)
        file_type = detect_file_type(input_path)
        signals.getFile.emit(fullpath,item,file_type)
        print(f"正在处理: {input_path} [类型: {file_type}]")
        handlers = {
            'DOCX': lambda: process_word(input_path, temp_dir),
            'DOC': lambda: process_word_legacy(input_path),
            'XLSX': lambda: process_excel(input_path),
            'XLS': lambda: process_excel_legacy(input_path),
            'PPTX': lambda: process_pptx(input_path),
            'PPT': lambda: process_ppt_legacy(input_path),
            'PDF': lambda: process_pdf(input_path, temp_dir),
            'Image': lambda: extract_text_from_image(input_path),
            'Text': lambda: open(input_path, 'r', encoding='utf-8', errors='ignore').read(),
            'Markdown': lambda: open(input_path, 'r', encoding='utf-8').read(),
            'ZIP': lambda: process_compressed_content(input_path,temp_dir),
            'RAR': lambda: process_compressed_content(input_path, temp_dir),
            'Unknown': lambda: process_unknown(input_path),
            'Empty': lambda: process_empty(input_path),
            'RTF': lambda: open(input_path, 'r', encoding='utf-8', errors='ignore').read(),
        }

        result = handlers.get(file_type, lambda: "")()

        # 处理压缩文件
        # if file_type in ['ZIP', 'RAR']:
        #     for extracted_file in result:
        #         preprocess_zip(extracted_file, temp_dir,ScanCallback, processed_files)
        #     return

        cleaned_content = clean_content(result)
        content_parts = split_content(cleaned_content)
        
        #将整个文本发送给model处理
        yield content_parts,fullpath
        # 写入CSV
#TODO 递归的处理对应zip 将文本归类于总的压缩文件上而不是压缩包内的文件
def preprocess_zip(input_path, writer, temp_dir,ScanCallback:Callable[[str,str,str,any],None], processed_files=None):
    if processed_files is None:
        processed_files = set()
    if os.path.isfile(input_path):
        file_id = f"{os.path.abspath(input_path)}-{os.path.getmtime(input_path)}"
        if file_id in processed_files:
            return
        processed_files.add(file_id)

        input_path = rename_file_based_on_type(input_path,temp_dir)
        file_type = detect_file_type(input_path)
        print(f"正在处理: {input_path} [类型: {file_type}]")
        ScanCallback(input_path,file_type)
        handlers = {
            'DOCX': lambda: process_word(input_path, temp_dir),
            'DOC': lambda: process_word_legacy(input_path),
            'XLSX': lambda: process_excel(input_path),
            'XLS': lambda: process_excel_legacy(input_path),
            'PPTX': lambda: process_pptx(input_path),
            'PPT': lambda: process_ppt_legacy(input_path),
            'PDF': lambda: process_pdf(input_path, temp_dir),
            'Image': lambda: extract_text_from_image(input_path),
            'Text': lambda: open(input_path, 'r', encoding='utf-8', errors='ignore').read(),
            'Markdown': lambda: open(input_path, 'r', encoding='utf-8').read(),
            'ZIP': lambda: process_compressed(input_path, os.path.join(temp_dir, "zip_extract")),
            'RAR': lambda: process_compressed(input_path, os.path.join(temp_dir, "rar_extract")),
            'Unknown': lambda: process_unknown(input_path),
            'Empty': lambda: process_empty(input_path),
            'RTF': lambda: open(input_path, 'r', encoding='utf-8', errors='ignore').read(),
        }

        result = handlers.get(file_type, lambda: "")()

        # 处理压缩文件
        if file_type in ['ZIP', 'RAR']:
            for extracted_file in result:
                preprocess_zip(extracted_file, writer, temp_dir,ScanCallback, processed_files)
            return

        # 处理文件名
        filename = os.path.basename(input_path)
        base_name = os.path.splitext(filename)[0]
        clean_name = re.sub(r'[^A-Za-z0-9]', '', base_name)

        # 处理内容
        cleaned_content = clean_content(result)
        content_parts = split_content(cleaned_content)
        ScanCallback(clean_name,file_type,content_parts)

    elif os.path.isdir(input_path):
        for item in os.listdir(input_path):
            preprocess_recursion(os.path.join(input_path, item), writer, temp_dir,ScanCallback, processed_files)

def start_preprocess(scandir:str,recursion:bool,signals):
    clear_com_cache()

    input_dir = scandir
    temp_dir = os.path.join(cfg.get(cfg.tempDir), "scan")
    if os.path.exists(temp_dir):
        # shutil.rmtree(temp_dir, ignore_errors=True)#在所有扫描结束后清除
        pass

    with OfficeSecurityConfig():
        # shutil.rmtree(temp_dir, ignore_errors=True)
        os.makedirs(temp_dir, exist_ok=True)
        if recursion:
            yield from preprocess_line(input_dir, temp_dir,signals)
        else:
            yield from preprocess_line(input_dir, temp_dir,signals)
